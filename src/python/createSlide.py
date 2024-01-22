import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(chapters):
    # Create a presentation object
    prs = Presentation()

    for chapter in chapters:
        # Split the chapter into paragraphs
        paragraphs = chapter.split('\n')

        # Process each paragraph
        for paragraph in paragraphs:
            words = paragraph.split()
            word_count = len(words)

            # Check if the paragraph fits the criteria for a new slide
            if word_count >= 20 and word_count <= 40 or (word_count < 20 and len(paragraphs) == 1):
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using a title and content layout
                title = slide.shapes.title
                title.text = "Chapter Summary"  # You can modify this to suit your title needs

                # Add the paragraph as a bullet point
                content = slide.placeholders[1]
                content.text = paragraph
                for paragraph in content.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(18)  # Adjust font size as needed
                        paragraph.alignment = PP_ALIGN.LEFT

            # If the paragraph is too short and it's not the only paragraph, accumulate words
            elif word_count < 20:
                accumulated_paragraph = paragraph
                while word_count < 30 and paragraphs:
                    next_paragraph = paragraphs.pop(0)
                    next_word_count = len(next_paragraph.split())
                    if word_count + next_word_count <= 30:
                        accumulated_paragraph += "\n" + next_paragraph
                        word_count += next_word_count
                    else:
                        paragraphs.insert(0, next_paragraph)
                        break

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = "Chapter Summary"

                content = slide.placeholders[1]
                content.text = accumulated_paragraph
                for paragraph in content.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        paragraph.alignment = PP_ALIGN.LEFT

    return prs

def process_markdown_file(md_file_path):
    # Read the markdown file
    with open(md_file_path, 'r', encoding='utf-8') as file:
        md_content = file.read()

    # Split the content into chapters based on headings (assuming '#' is used for chapter headings)
    chapters = re.split(r'\n#+ ', md_content)

    # Remove empty strings and strip leading/trailing whitespaces
    chapters = [chapter.strip() for chapter in chapters if chapter.strip()]

    return chapters

# Example Markdown file path
md_file_path = 'pirahansiah/pirahansiah.github.io/src/books/pkm/chapters/chapter02.md'  # Replace with actual file path

# Process the Markdown file to get chapters
chapters = process_markdown_file(md_file_path)

# Create the presentation with the obtained chapters
presentation = create_presentation(chapters)

# Save the presentation
presentation_file_path = 'pirahansiah/pirahansiah.github.io/src/books/pkm/chapters/presentation_from_md.pptx'
presentation.save(presentation_file_path)

presentation_file_path

